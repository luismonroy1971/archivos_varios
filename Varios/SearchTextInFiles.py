
import os
import PyPDF2
import openpyxl
from pptx import Presentation
from docx import Document

def search_in_pdf(file_path, text):
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            if text.lower() in page.extract_text().lower():
                return True
    return False

def search_in_excel(file_path, text):
    workbook = openpyxl.load_workbook(file_path)
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows(values_only=True):
            for cell in row:
                if cell and text.lower() in str(cell).lower():
                    return True
    return False

def search_in_ppt(file_path, text):
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and text.lower() in shape.text.lower():
                return True
    return False

def search_in_doc(file_path, text):
    doc = Document(file_path)
    for paragraph in doc.paragraphs:
        if text.lower() in paragraph.text.lower():
            return True
    return False

def search_text_in_files(folder_path, text):
    for root, _, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            try:
                if file.endswith('.pdf') and search_in_pdf(file_path, text):
                    print(f"Found in: {file_path}")
                elif file.endswith('.xlsx') and search_in_excel(file_path, text):
                    print(f"Found in: {file_path}")
                elif file.endswith('.pptx') and search_in_ppt(file_path, text):
                    print(f"Found in: {file_path}")
                elif file.endswith('.docx') and search_in_doc(file_path, text):
                    print(f"Found in: {file_path}")
            except Exception as e:
                print(f"Error processing {file_path}: {e}")

# Example usage:
# search_text_in_files('path_to_folder', 'target_text')
